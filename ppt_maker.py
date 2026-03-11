import os, json, requests, re, time
from urllib.parse import quote_plus
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from io import BytesIO
from PIL import Image

# ==========================================
# CONFIGURATION
# ==========================================
API_KEY = os.getenv('GEMINI_API_KEY')
client = genai.Client(api_key=API_KEY)

# Professional Color Palette
DARK_BG       = RGBColor(3, 0, 46)       # Deep navy-black
ACCENT_BG     = RGBColor(30, 60, 114)       # Section header blue
CARD_BG       = RGBColor(34, 40, 60)        # Slightly lighter card
WHITE         = RGBColor(255, 255, 255)
LIGHT_GRAY    = RGBColor(200, 210, 225)
ACCENT_GOLD   = RGBColor(255, 196, 61)      # Gold accent
ACCENT_CYAN   = RGBColor(0, 188, 212)       # Teal accent for highlights
SUBTLE_GRAY   = RGBColor(120, 130, 150)
BULLET_COLOR  = RGBColor(0, 188, 212)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"

# ==========================================
# GEMINI HELPERS
# ==========================================
GEMINI_MODELS = [
    "gemini-2.5-flash-lite",
    "gemini-2.0-flash",
    "gemini-2.0-flash-lite",
    "gemini-1.5-flash",
    "gemini-1.5-flash-8b",
]

def parse_json_response(text):
    cleaned = re.sub(r'^```(?:json)?\s*\n?', '', text.strip(), flags=re.MULTILINE)
    cleaned = re.sub(r'\n?```\s*$', '', cleaned.strip(), flags=re.MULTILINE)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass
    match = re.search(r'\[.*\]', cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass
    raise ValueError(f"Could not parse JSON from response:\n{text[:500]}")

def gemini_generate(prompt, max_retries=3):
    """Try each model in GEMINI_MODELS; on rate limit move to the next."""
    errors = []
    for model in GEMINI_MODELS:
        for attempt in range(max_retries):
            try:
                return client.models.generate_content(
                    model=model,
                    contents=prompt,
                    config=types.GenerateContentConfig(response_mime_type='application/json')
                )
            except Exception as e:
                err = str(e)
                if '429' in err or 'RESOURCE_EXHAUSTED' in err:
                    if attempt < max_retries - 1:
                        wait = 15 * (attempt + 1)
                        print(f"  ⏳ {model} rate limited. Waiting {wait}s... ({attempt+1}/{max_retries})")
                        time.sleep(wait)
                        continue
                    print(f"  ⚠️ {model} exhausted, trying next model...")
                    errors.append(f"{model}: rate limited")
                    break
                else:
                    raise
    raise RuntimeError(f"All Gemini models exhausted. Tried: {', '.join(errors)}")

# ==========================================
# IMAGE SCRAPER (Bing) — only used for diagram slides
# ==========================================
def scrape_image(query):
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.5",
    }
    try:
        url = f"https://www.bing.com/images/search?q={quote_plus(query)}&first=1"
        resp = requests.get(url, headers=headers, timeout=15)
        if resp.status_code != 200:
            return None, f"HTTP {resp.status_code}"
        img_urls = re.findall(r'murl&quot;:&quot;(https?://[^&]+?)&quot;', resp.text)
        if not img_urls:
            img_urls = re.findall(r'src2?="(https?://[^"]+\.(?:jpg|jpeg|png|webp))', resp.text)
        for img_url in img_urls[:5]:
            try:
                r = requests.get(img_url, headers=headers, timeout=10)
                if r.status_code == 200 and len(r.content) > 2000:
                    stream = BytesIO(r.content)
                    img = Image.open(stream)
                    if img.format not in ('BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF'):
                        buf = BytesIO()
                        img.convert('RGBA').save(buf, format='PNG')
                        buf.seek(0)
                        return buf, "OK"
                    stream.seek(0)
                    return stream, "OK"
            except Exception:
                continue
        return None, "No downloadable image"
    except Exception as e:
        return None, str(e)

# ==========================================
# SHAPE HELPERS
# ==========================================
def set_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_char="▸"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    shape = slide.shapes.add_shape(
        1, left, top, width, Pt(3)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        5, left, top, width, height  # 5 = rounded rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ==========================================
# LOGO HELPER
# ==========================================
LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo.png")

def _lock_shape(shape):
    """Lock a shape so it can't be selected, moved, resized or deleted in PowerPoint."""
    cNvPr = shape._element.find(qn('p:cNvPr'))
    if cNvPr is None:
        # For picture shapes the path is different
        cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
    # Add locks to prevent editing
    cNvPicPr = shape._element.find('.//' + qn('p:cNvPicPr'))
    if cNvPicPr is not None:
        locks = cNvPicPr.find(qn('a:picLocks'))
        if locks is None:
            from lxml import etree
            locks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
        locks.set('noSelect', '1')
        locks.set('noMove', '1')
        locks.set('noResize', '1')
        locks.set('noRot', '1')
        locks.set('noChangeAspect', '1')

def add_logo(slide, left=None, top=None, height=Inches(0.5)):
    """Add locked branding logo to a slide (top-left by default)."""
    if not os.path.exists(LOGO_PATH):
        return
    if left is None:
        left = Inches(0.3)
    if top is None:
        top = Inches(0.2)
    pic = slide.shapes.add_picture(LOGO_PATH, left, top, height=height)
    _lock_shape(pic)

# ==========================================
# SLIDE BUILDERS
# ==========================================
def build_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Accent line top
    add_accent_line(slide, Inches(1.5), Inches(2.5), Inches(10.33))

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.33), Inches(1.5),
                title, font_size=42, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.33), Inches(1),
                subtitle, font_size=20, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Accent line bottom
    add_accent_line(slide, Inches(1.5), Inches(5.2), Inches(10.33))

    add_logo(slide)

def build_section_slide(prs, section_number, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Section number
    add_textbox(slide, Inches(1), Inches(2.3), Inches(11.33), Inches(1),
                f"SECTION {section_number}", font_size=18, color=ACCENT_CYAN,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Section title
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11.33), Inches(1.5),
                section_title, font_size=38, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_accent_line(slide, Inches(5), Inches(4.7), Inches(3.33), ACCENT_GOLD)

    add_logo(slide)

def build_content_slide(prs, slide_num, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8),
                title, font_size=30, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Accent underline
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(11.73))

    # Content card
    add_rounded_card(slide, Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.5))

    # Bullet list inside card
    if isinstance(bullets, list):
        add_bullet_list(slide, Inches(1.3), Inches(1.8), Inches(10.73), Inches(5),
                        bullets, font_size=18, color=WHITE)
    else:
        # If body is a string, split into lines
        lines = [l.strip('- •▸') .strip() for l in str(bullets).split('\n') if l.strip()]
        if len(lines) <= 1:
            lines = [s.strip() for s in str(bullets).split('.') if s.strip()]
        add_bullet_list(slide, Inches(1.3), Inches(1.8), Inches(10.73), Inches(5),
                        lines[:10], font_size=18, color=WHITE)

    add_logo(slide)

def build_diagram_slide(prs, slide_num, title, bullets, search_query):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8),
                title, font_size=28, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(11.73))

    # LEFT: content card (narrower)
    add_rounded_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))

    if isinstance(bullets, list):
        items = bullets
    else:
        items = [l.strip('- •▸').strip() for l in str(bullets).split('\n') if l.strip()]
        if len(items) <= 1:
            items = [s.strip() for s in str(bullets).split('.') if s.strip()]

    add_bullet_list(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(5),
                    items[:8], font_size=18, color=WHITE)

    # RIGHT: diagram image
    img_stream, status = scrape_image(search_query)
    if img_stream:
        print(f"  Slide {slide_num}: ✅ Diagram fetched")
        # Image card background
        add_rounded_card(slide, Inches(6.8), Inches(1.5), Inches(6.03), Inches(5.5),
                         fill_color=RGBColor(25, 30, 48))
        slide.shapes.add_picture(img_stream, Inches(7.0), Inches(1.7),
                                 width=Inches(5.6), height=Inches(5.0))
    else:
        print(f"  Slide {slide_num}: ❌ No diagram | {status}")
        # Placeholder card
        add_rounded_card(slide, Inches(6.8), Inches(1.5), Inches(6.03), Inches(5.5),
                         fill_color=RGBColor(25, 30, 48))
        add_textbox(slide, Inches(7.0), Inches(3.5), Inches(5.5), Inches(1),
                    "[Diagram: " + search_query + "]",
                    font_size=16, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)

    add_logo(slide)

def build_summary_slide(prs, title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8),
                title, font_size=32, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(11.73), ACCENT_GOLD)

    # Summary points in two columns
    mid = (len(points) + 1) // 2
    left_items = points[:mid]
    right_items = points[mid:]

    add_rounded_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    add_bullet_list(slide, Inches(1.2), Inches(1.9), Inches(4.9), Inches(4.8),
                    left_items, font_size=18, color=WHITE, bullet_char="✦")

    if right_items:
        add_rounded_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.3))
        add_bullet_list(slide, Inches(7.3), Inches(1.9), Inches(4.9), Inches(4.8),
                        right_items, font_size=18, color=WHITE, bullet_char="✦")

    add_logo(slide)

def build_thank_you_slide(prs, class_num, subject, chapter):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
                "Thank You", font_size=48, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_accent_line(slide, Inches(5), Inches(4.2), Inches(3.33), ACCENT_GOLD)

    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(1),
                f"Class {class_num} {subject} — {chapter}",
                font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_logo(slide)

# ==========================================
# MAIN — PROMPT + ASSEMBLY
# ==========================================
GEMINI_PROMPT_TEMPLATE = """You are an expert {subject} teacher. Create a structured presentation on "{chapter}" for Class {class_num} students.

Return a JSON array of exactly 30 slide objects. Each slide must have these keys:
- "type": one of "title", "section", "content", "diagram", "summary"
- "title": slide heading
- "body": array of 4-6 concise bullet point strings (not needed for title/section)
- "search_query": a short image search query ONLY for "diagram" type slides (e.g. "labeled diagram of ..."). Leave empty string for other types.
- "section_number": integer, only for "section" type slides

STRUCTURE RULES:
- Slide 1: type "title" (main title slide)
- Then organize into 4-6 logical sections based on the chapter content
- Each section starts with a "section" slide
- Within each section: mix of "content" (text-only explanation) and "diagram" slides (where a labeled diagram is essential)
- Use "diagram" type ONLY for slides that genuinely need a visual. About 8-10 diagram slides total.
- Last slide: type "summary"
- Keep bullet points concise (max 15 words each)

Return ONLY the JSON array, no other text."""

GEMINI_PROMPT_TEXT_ONLY = """You are an expert {subject} teacher. Create a detailed, text-heavy presentation on "{chapter}" for Class {class_num} students.

Return a JSON array of exactly 30 slide objects. Each slide must have these keys:
- "type": one of "title", "section", "content", "summary"
- "title": slide heading
- "body": array of 6-8 detailed, informative bullet point strings. Each bullet should be a complete explanation (15-25 words) that teaches the concept clearly. NOT needed for title/section slides.
- "section_number": integer, only for "section" type slides

STRUCTURE RULES:
- Slide 1: type "title" (main title slide, body should have a one-line subtitle)
- Then organize into 4-6 logical sections based on the chapter content
- Each section starts with a "section" slide
- All other slides are "content" type — NO diagram slides. Every slide must have detailed text.
- Last slide: type "summary" with key revision points
- Make bullet points detailed and educational — students should be able to learn from reading the slides alone
- Cover definitions, processes, examples, comparisons, and important facts
- Include exam-relevant points, mnemonics, and key differences where applicable

Return ONLY the JSON array, no other text."""

def create_ppt():
    # --- User Input ---
    print("="*50)
    print("   📚 AI Presentation Maker")
    print("="*50)
    class_num = input("  Enter Class (e.g. 10): ").strip()
    subject   = input("  Enter Subject (e.g. Biology): ").strip()
    chapter   = input("  Enter Chapter (e.g. Human Life Processes): ").strip()
    use_images = input("  Include images/diagrams? (yes/no): ").strip().lower() in ('yes', 'y')
    print("="*50)

    if use_images:
        prompt = GEMINI_PROMPT_TEMPLATE.format(
            class_num=class_num, subject=subject, chapter=chapter
        )
    else:
        prompt = GEMINI_PROMPT_TEXT_ONLY.format(
            class_num=class_num, subject=subject, chapter=chapter
        )

    print(f"\n🤖 Generating 30-slide structure for Class {class_num} {subject}: {chapter}...")
    response = gemini_generate(prompt)
    slides_data = parse_json_response(response.text)
    print(f"  ✅ Got {len(slides_data)} slides from Gemini\n")

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    slide_num = 0
    print("--- Building Slides ---")

    for data in slides_data:
        slide_type = data.get('type', 'content')
        title = data.get('title', '')
        body = data.get('body', [])
        # Ensure body is never empty for content/diagram/summary slides
        if slide_type in ('content', 'diagram', 'summary') and not body:
            body = [f"Key points about {title}"]
        query = data.get('search_query', '')
        sec_num = data.get('section_number', '')

        if slide_type == 'title':
            subtitle = body[0] if isinstance(body, list) and body else f"Class {class_num} {subject}"
            build_title_slide(prs, title, subtitle)
            print(f"  📘 Title Slide: {title}")

        elif slide_type == 'section':
            build_section_slide(prs, sec_num, title)
            print(f"  📂 Section {sec_num}: {title}")

        elif slide_type == 'diagram':
            slide_num += 1
            if use_images:
                build_diagram_slide(prs, slide_num, title, body, query)
                print(f"  🖼️  Slide {slide_num}: {title} [DIAGRAM]")
            else:
                build_content_slide(prs, slide_num, title, body)
                print(f"  📝 Slide {slide_num}: {title} [TEXT ONLY]")

        elif slide_type == 'summary':
            items = body if isinstance(body, list) else [body]
            build_summary_slide(prs, title, items)
            print(f"  📋 Summary: {title}")

        else:  # content
            slide_num += 1
            build_content_slide(prs, slide_num, title, body)
            print(f"  📝 Slide {slide_num}: {title}")

    # Final thank you
    build_thank_you_slide(prs, class_num, subject, chapter)
    print("  🎉 Thank You slide added")

    # Safe filename from user input
    safe_name = re.sub(r'[^\w\s-]', '', chapter).strip().replace(' ', '_')
    filename = f"Class{class_num}_{subject}_{safe_name}.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation saved: {filename}")
    os.startfile(filename)

if __name__ == "__main__":
    create_ppt()