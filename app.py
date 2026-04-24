import streamlit as st
import os, json, requests, re, time, zipfile
import ast
from urllib.parse import quote_plus
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from io import BytesIO
from PIL import Image
from fpdf import FPDF
import tempfile

# ==========================================
# CONFIGURATION
# ==========================================
API_KEY = os.getenv('GEMINI_API_KEY')
client = genai.Client(api_key=API_KEY)

# ── Dual-theme palette ──────────────────────────────────────────────
HEADING_BG    = RGBColor(10,  25,  80)
HEADING_TEXT  = RGBColor(255, 196, 61)
HEADING_SUB   = RGBColor(220, 235, 255)

CONTENT_BG    = RGBColor(255, 255, 255)
CONTENT_TITLE = RGBColor(10,  25,  80)
CONTENT_BODY  = RGBColor(30,  30,  30)
CONTENT_BULLET= RGBColor(10,  25,  80)
CARD_BG       = RGBColor(240, 244, 255)

ACCENT_GOLD   = RGBColor(255, 196, 61)
ACCENT_BLUE   = RGBColor(10,  25,  80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo.png")

# ==========================================
# GEMINI HELPERS
# ==========================================
GEMINI_MODELS = [
    "gemini-2.5-flash",
    "gemini-2.5-flash-lite",
    "gemini-2.0-flash",
    "gemini-2.0-flash-lite",
    "gemini-1.5-flash",
    "gemini-1.5-flash-8b",
]

def parse_json_response(text):
    def _strip_fences(raw):
        cleaned_local = re.sub(r'^```(?:json)?\s*\n?', '', raw.strip(), flags=re.MULTILINE)
        return re.sub(r'\n?```\s*$', '', cleaned_local.strip(), flags=re.MULTILINE)

    def _extract_first_json_array(raw):
        start = raw.find('[')
        if start == -1:
            return None
        depth, in_str, escape = 0, False, False
        for i in range(start, len(raw)):
            ch = raw[i]
            if in_str:
                if escape: escape = False
                elif ch == '\\': escape = True
                elif ch == '"': in_str = False
                continue
            if ch == '"': in_str = True
            elif ch == '[': depth += 1
            elif ch == ']':
                depth -= 1
                if depth == 0:
                    return raw[start:i + 1]
        return None

    def _sanitize_json_like(raw):
        s = raw
        s = s.replace('\u201c', '"').replace('\u201d', '"')
        s = s.replace('\u2018', "'").replace('\u2019', "'")
        s = re.sub(r',\s*([}\]])', r'\1', s)
        s = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', s)
        return s.strip()

    cleaned = _strip_fences(text)
    candidates = [cleaned]
    first_array = _extract_first_json_array(cleaned)
    if first_array:
        candidates.append(first_array)
    candidates.append(_sanitize_json_like(cleaned))
    if first_array:
        candidates.append(_sanitize_json_like(first_array))

    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, list):
                return parsed
        except json.JSONDecodeError:
            continue

    for candidate in candidates:
        try:
            parsed = ast.literal_eval(candidate)
            if isinstance(parsed, list):
                return parsed
        except Exception:
            continue

    raise ValueError("Could not parse JSON")

def gemini_generate(prompt, response_mime='application/json'):
    for model in GEMINI_MODELS:
        try:
            return client.models.generate_content(
                model=model,
                contents=prompt,
                config=types.GenerateContentConfig(response_mime_type=response_mime)
            )
        except:
            continue
    raise RuntimeError("All Gemini models failed")

def gemini_generate_text(prompt):
    for model in GEMINI_MODELS:
        try:
            return client.models.generate_content(model=model, contents=prompt).text
        except:
            continue
    raise RuntimeError("All Gemini models failed")

# ==========================================
# IMAGE FETCH
# ==========================================
def scrape_image(query):
    try:
        url = f"https://www.bing.com/images/search?q={quote_plus(query)}"
        resp = requests.get(url, timeout=10)
        img_urls = re.findall(r'murl&quot;:&quot;(https?://[^&]+?)&quot;', resp.text)
        for u in img_urls[:3]:
            try:
                r = requests.get(u, timeout=5)
                if r.status_code == 200:
                    return BytesIO(r.content)
            except:
                continue
    except:
        pass
    return None

# ==========================================
# PPT GENERATION (UNCHANGED CORE)
# ==========================================
def generate_ppt(class_num, subject, chapter):
    prompt = f"Create 20 slides JSON for {chapter}"
    response = gemini_generate(prompt)
    slides = parse_json_response(response.text)

    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = "\n".join(s.get("body", []))
        slide.placeholders[1].text = body

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf, f"{chapter}.pptx"

# ==========================================
# PDF NOTES (UNCHANGED CORE)
# ==========================================
def generate_notes(class_num, subject, chapter):
    text = gemini_generate_text(f"Write notes for {chapter}")

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for line in text.split("\n"):
        pdf.multi_cell(0, 6, line)

    buf = BytesIO()
    buf.write(pdf.output())
    buf.seek(0)
    return buf, f"{chapter}.pdf"

# ==========================================
# STREAMLIT UI
# ==========================================
st.title("📚 AI Education Generator")

class_num = st.text_input("Class")
subject = st.text_input("Subject")
chapter = st.text_input("Chapter")

generate = st.button("Generate")

if generate:
    ppt_buf, ppt_name = generate_ppt(class_num, subject, chapter)
    notes_buf, notes_name = generate_notes(class_num, subject, chapter)

    zip_buf = BytesIO()
    with zipfile.ZipFile(zip_buf, 'w') as z:
        z.writestr(ppt_name, ppt_buf.getvalue())
        z.writestr(notes_name, notes_buf.getvalue())

    zip_buf.seek(0)

    st.download_button(
        "Download ZIP",
        zip_buf,
        file_name=f"{chapter}.zip"
    )
